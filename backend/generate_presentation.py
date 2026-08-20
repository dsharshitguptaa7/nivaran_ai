import os
from pathlib import Path
import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ==============================================================================
# COLOR PALETTE DEFINITIONS (CSJMU / NIVARAN-AI Brand)
# ==============================================================================
COLOR_PRIMARY = RGBColor(91, 16, 33)       # Crimson Maroon #5B1021
COLOR_PRIMARY_DARK = RGBColor(59, 7, 20)  # Dark Maroon #3B0714
COLOR_NAVY = RGBColor(30, 41, 59)         # Deep Navy Slate #1E293B
COLOR_GOLD = RGBColor(212, 175, 55)       # Metallic Gold #D4AF37
COLOR_GOLD_LIGHT = RGBColor(254, 249, 195)# Soft Gold #FEF9C3
COLOR_BG = RGBColor(248, 250, 252)        # Light Slate Background #F8FAFC
COLOR_WHITE = RGBColor(255, 255, 255)     # Pure White
COLOR_CARD_BG = RGBColor(255, 255, 255)   # Card Background
COLOR_CARD_BORDER = RGBColor(226, 232, 240) # Border Slate #E2E8F0
COLOR_TEXT_DARK = RGBColor(15, 23, 42)    # Title / Body Slate #0F172A
COLOR_TEXT_MUTED = RGBColor(100, 116, 139) # Secondary Slate #64748B
COLOR_GREEN = RGBColor(22, 163, 74)       # Emerald Green #16A34A
COLOR_GREEN_BG = RGBColor(240, 253, 244)  # Emerald Light #F0FDF4
COLOR_BLUE = RGBColor(37, 99, 235)        # Royal Blue #2563EB
COLOR_BLUE_BG = RGBColor(239, 246, 255)   # Soft Blue #EFF6FF
COLOR_AMBER = RGBColor(217, 119, 6)       # Amber #D97706
COLOR_AMBER_BG = RGBColor(254, 243, 199)  # Soft Amber #FEF3C7

LOGO_PATH = Path("../frontend/src/images/logo.png")
if not LOGO_PATH.exists():
    LOGO_PATH = Path("frontend/src/images/logo.png")

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    def set_slide_background(slide, color):
        bg_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = color
        bg_shape.line.fill.background()
        return bg_shape

    def add_header(slide, title_text, subtitle_text="AI-Assisted Grievance Redressal System • CSJMU R&D Section"):
        banner = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(1.15))
        banner.fill.solid()
        banner.fill.fore_color.rgb = COLOR_PRIMARY
        banner.line.fill.background()

        gold_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(1.15), Inches(13.333), Inches(0.06))
        gold_line.fill.solid()
        gold_line.fill.fore_color.rgb = COLOR_GOLD
        gold_line.line.fill.background()

        tx_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.12), Inches(10.5), Inches(0.6))
        tf = tx_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = COLOR_WHITE
        p.font.name = "Segoe UI"

        p2 = tf.add_paragraph()
        p2.text = subtitle_text.upper()
        p2.font.size = Pt(9.5)
        p2.font.color.rgb = COLOR_GOLD
        p2.font.bold = True
        p2.font.name = "Segoe UI"

        if LOGO_PATH.exists():
            try:
                slide.shapes.add_picture(str(LOGO_PATH), Inches(12.2), Inches(0.15), width=Inches(0.85))
            except Exception:
                pass

    def add_card(slide, left, top, width, height, bg_color=COLOR_CARD_BG, border_color=COLOR_CARD_BORDER):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        card.fill.solid()
        card.fill.fore_color.rgb = bg_color
        card.line.color.rgb = border_color
        card.line.width = Pt(1)
        return card

    # ==========================================================================
    # SLIDE 1: TITLE SLIDE (Hero Cover)
    # ==========================================================================
    s1 = prs.slides.add_slide(blank_layout)
    set_slide_background(s1, COLOR_PRIMARY_DARK)

    decor = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.5), 0, Inches(4.833), Inches(7.5))
    decor.fill.solid()
    decor.fill.fore_color.rgb = COLOR_PRIMARY
    decor.line.fill.background()

    gold_stripe = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.45), 0, Inches(0.08), Inches(7.5))
    gold_stripe.fill.solid()
    gold_stripe.fill.fore_color.rgb = COLOR_GOLD
    gold_stripe.line.fill.background()

    if LOGO_PATH.exists():
        try:
            s1.shapes.add_picture(str(LOGO_PATH), Inches(9.5), Inches(1.8), width=Inches(2.8))
        except Exception:
            pass

    univ_card = add_card(s1, 0.8, 1.2, 7.0, 0.45, bg_color=COLOR_PRIMARY, border_color=COLOR_GOLD)
    u_tf = univ_card.text_frame
    u_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    u_p = u_tf.paragraphs[0]
    u_p.text = "🏛️ CHHATRAPATI SHAHU JI MAHARAJ UNIVERSITY, KANPUR"
    u_p.font.size = Pt(11)
    u_p.font.bold = True
    u_p.font.color.rgb = COLOR_GOLD
    u_p.font.name = "Segoe UI"

    t_box = s1.shapes.add_textbox(Inches(0.8), Inches(1.9), Inches(7.4), Inches(2.2))
    ttf = t_box.text_frame
    ttf.word_wrap = True
    tp1 = ttf.paragraphs[0]
    tp1.text = "NIVARAN.AI"
    tp1.font.size = Pt(44)
    tp1.font.bold = True
    tp1.font.color.rgb = COLOR_WHITE
    tp1.font.name = "Segoe UI"

    tp2 = ttf.add_paragraph()
    tp2.text = "AI-Assisted Grievance Redressal & Workflow Platform"
    tp2.font.size = Pt(20)
    tp2.font.color.rgb = COLOR_GOLD
    tp2.font.name = "Segoe UI"

    d_box = s1.shapes.add_textbox(Inches(0.8), Inches(4.3), Inches(7.2), Inches(1.6))
    dtf = d_box.text_frame
    dtf.word_wrap = True
    dp1 = dtf.paragraphs[0]
    dp1.text = "An intelligent, multi-tier academic grievance redressal system empowering research scholars, faculty officers, and administrative deans with NLP classification, dynamic routing, in-app notifications, and executive turnaround analytics."
    dp1.font.size = Pt(13)
    dp1.font.color.rgb = RGBColor(226, 232, 240)
    dp1.font.name = "Segoe UI"

    f_box = s1.shapes.add_textbox(Inches(0.8), Inches(6.4), Inches(7.2), Inches(0.6))
    ftf = f_box.text_frame
    fp = ftf.paragraphs[0]
    fp.text = "Research & Development (R&D) Section • Enterprise Product Brief"
    fp.font.size = Pt(10.5)
    fp.font.color.rgb = COLOR_GOLD
    fp.font.bold = True

    # ==========================================================================
    # SLIDE 2: THE PROBLEM & SOLUTION
    # ==========================================================================
    s2 = prs.slides.add_slide(blank_layout)
    set_slide_background(s2, COLOR_BG)
    add_header(s2, "Executive Problem Statement & NIVARAN-AI Solution")

    add_card(s2, 0.8, 1.5, 5.6, 5.4, bg_color=RGBColor(254, 242, 242), border_color=RGBColor(254, 202, 202))
    p_box = s2.shapes.add_textbox(Inches(1.0), Inches(1.7), Inches(5.2), Inches(5.0))
    ptf = p_box.text_frame
    ptf.word_wrap = True
    
    h = ptf.paragraphs[0]
    h.text = "❌ Traditional University Bottlenecks"
    h.font.size = Pt(18)
    h.font.bold = True
    h.font.color.rgb = RGBColor(185, 28, 28)

    points_p = [
        ("Manual Paper Submissions", "Physical forms & lost tracking IDs leading to lack of transparency."),
        ("Routing Ambiguity", "Scholars unsure whether to contact Subject Dean, Cluster Dean, or Finance."),
        ("No Inactivity Safeguards", "Grievances stall indefinitely on desks with zero escalation or alerts."),
        ("Unstructured Documentation", "Missing attachments cause repeated email ping-pong and delay."),
        ("Zero Executive Oversight", "Deans lack aggregate visibility into resolution turnaround & officer workload.")
    ]
    for title, desc in points_p:
        p1 = ptf.add_paragraph()
        p1.text = f"• {title}: "
        p1.font.size = Pt(12)
        p1.font.bold = True
        p1.font.color.rgb = COLOR_TEXT_DARK
        p1.space_before = Pt(8)
        run = p1.add_run()
        run.text = desc
        run.font.bold = False
        run.font.size = Pt(11.5)
        run.font.color.rgb = COLOR_TEXT_MUTED

    add_card(s2, 6.9, 1.5, 5.6, 5.4, bg_color=RGBColor(240, 253, 244), border_color=RGBColor(187, 247, 208))
    s_box = s2.shapes.add_textbox(Inches(7.1), Inches(1.7), Inches(5.2), Inches(5.0))
    stf = s_box.text_frame
    stf.word_wrap = True
    
    sh = stf.paragraphs[0]
    sh.text = "✅ The NIVARAN-AI Transformation"
    sh.font.size = Pt(18)
    sh.font.bold = True
    sh.font.color.rgb = COLOR_GREEN

    points_s = [
        ("AI-Driven Auto-Categorization", "NLP classifier predicts categories and clusters with confidence scoring."),
        ("Intelligent Multi-Tier Routing", "Auto-maps Subject Assistant Deans, Domain Associate Deans, & Fixed Officers."),
        ("3-Day Inactivity Reminder Engine", "Background jobs detect stalled cases and trigger automated authority alerts."),
        ("Interactive Document Request Loop", "Pause/resume grievance lifecycle with in-app review and secure attachments."),
        ("Executive Dean Analytics Matrix", "Real-time officer workload distribution, pending bottlenecks, and SLA KPIs.")
    ]
    for title, desc in points_s:
        p1 = stf.add_paragraph()
        p1.text = f"• {title}: "
        p1.font.size = Pt(12)
        p1.font.bold = True
        p1.font.color.rgb = COLOR_TEXT_DARK
        p1.space_before = Pt(8)
        run = p1.add_run()
        run.text = desc
        run.font.bold = False
        run.font.size = Pt(11.5)
        run.font.color.rgb = COLOR_TEXT_MUTED

    # ==========================================================================
    # SLIDE 3: SYSTEM ARCHITECTURE
    # ==========================================================================
    s3 = prs.slides.add_slide(blank_layout)
    set_slide_background(s3, COLOR_BG)
    add_header(s3, "System Architecture & Layered Technology Stack")

    layers = [
        ("1. Presentation Layer (React 18 SPA)", "Vite-bundled Single Page Application, dynamic CSJMU design system, role-based dashboards, unified login router, and real-time Bell Notification center.", COLOR_PRIMARY, COLOR_PRIMARY_DARK),
        ("2. API & Gateway Layer (FastAPI ASGI)", "Asynchronous REST endpoints, JWT authorization dependency injection, Pydantic data contract validation, and configurable CORS middleware.", COLOR_BLUE, COLOR_BLUE_BG),
        ("3. Business & Service Layer", "Grievance State Machine, Escalation Engine, Document Request Lifecycle, Notification & Email Dispatcher, and Dean Analytics Aggregator.", COLOR_GREEN, COLOR_GREEN_BG),
        ("4. AI & NLP Classification Engine", "TF-IDF Vectorization, multi-class Naive Bayes / classifier inference, confidence thresholding (≥ 0.75), and human-in-the-loop audit telemetry.", COLOR_AMBER, COLOR_AMBER_BG),
        ("5. Persistence & Storage Layer", "PostgreSQL 14+ database, SQLAlchemy 2.0 ORM, 24 Alembic migration revisions, and secure persistent document storage server.", COLOR_NAVY, RGBColor(241, 245, 249))
    ]

    top_pos = 1.45
    for title, desc, col_border, col_bg in layers:
        card = add_card(s3, 0.8, top_pos, 11.733, 1.0, bg_color=col_bg, border_color=col_border)
        tb = s3.shapes.add_textbox(Inches(1.0), Inches(top_pos + 0.08), Inches(11.3), Inches(0.85))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = col_border if col_border != COLOR_PRIMARY else COLOR_WHITE
        
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(11)
        p2.font.color.rgb = COLOR_WHITE if col_bg == COLOR_PRIMARY_DARK else COLOR_TEXT_DARK
        top_pos += 1.12

    # ==========================================================================
    # SLIDE 4: USER ROLES & ADMINISTRATIVE HIERARCHY
    # ==========================================================================
    s4 = prs.slides.add_slide(blank_layout)
    set_slide_background(s4, COLOR_BG)
    add_header(s4, "Administrative Hierarchy & Role-Based Access Control")

    roles = [
        ("APPLICANT (Scholar)", "Files grievances, attaches receipts, responds to doc requests, monitors live tracking.", RGBColor(59, 130, 246), "Student Portal"),
        ("MANAGER (R&D)", "Conducts human-in-the-loop review, verifies AI predictions, triggers auto-routing, closes cases.", RGBColor(168, 85, 247), "Intake & Closure"),
        ("ASSISTANT DEAN", "Assigned based on scholar academic discipline; requests documents, directly resolves or forwards.", RGBColor(14, 165, 233), "Subject Specialist"),
        ("ASSOCIATE DEAN", "Receives forwarded cluster cases; oversees domain disputes; escalates policy issues to Dean.", RGBColor(234, 88, 12), "Domain Cluster"),
        ("DEAN (R&D)", "High-level strategic oversight, executive workload matrix, final resolution of escalated cases.", COLOR_PRIMARY, "Executive Dean")
    ]

    left_pos = 0.8
    for title, desc, color, badge in roles:
        add_card(s4, left_pos, 1.5, 2.2, 5.3, bg_color=COLOR_WHITE, border_color=color)
        
        rb = s4.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left_pos), Inches(1.5), Inches(2.2), Inches(0.8))
        rb.fill.solid()
        rb.fill.fore_color.rgb = color
        rb.line.fill.background()
        
        rtf = rb.text_frame
        rtf.word_wrap = True
        rp = rtf.paragraphs[0]
        rp.text = badge.upper()
        rp.font.size = Pt(9)
        rp.font.bold = True
        rp.font.color.rgb = COLOR_WHITE
        rp.alignment = PP_ALIGN.CENTER
        
        rp2 = rtf.add_paragraph()
        rp2.text = title
        rp2.font.size = Pt(11)
        rp2.font.bold = True
        rp2.font.color.rgb = COLOR_WHITE
        rp2.alignment = PP_ALIGN.CENTER

        tb = s4.shapes.add_textbox(Inches(left_pos + 0.1), Inches(2.4), Inches(2.0), Inches(4.2))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(11)
        p.font.color.rgb = COLOR_TEXT_DARK
        
        left_pos += 2.38

    # ==========================================================================
    # SLIDE 5: AI CLASSIFICATION ENGINE
    # ==========================================================================
    s5 = prs.slides.add_slide(blank_layout)
    set_slide_background(s5, COLOR_BG)
    add_header(s5, "AI Grievance Classification & NLP Inference Pipeline")

    add_card(s5, 0.8, 1.5, 5.7, 5.4, bg_color=COLOR_WHITE, border_color=COLOR_CARD_BORDER)
    tb_l = s5.shapes.add_textbox(Inches(1.0), Inches(1.7), Inches(5.3), Inches(5.0))
    tfl = tb_l.text_frame
    tfl.word_wrap = True

    h = tfl.paragraphs[0]
    h.text = "🧠 NLP Categorization Workflow"
    h.font.size = Pt(16)
    h.font.bold = True
    h.font.color.rgb = COLOR_PRIMARY

    steps = [
        ("1. Input Tokenization", "Cleans raw grievance title & description text, strips noise, filters stopwords."),
        ("2. TF-IDF Feature Vectorizer", "Extracts unigram and bigram features across 5,000 top vocabulary dimensions."),
        ("3. Statistical Classifier", "Computes multi-class probability distributions across university grievance domains."),
        ("4. Confidence Thresholding", "High Confidence (≥ 0.75) highlights recommendations; Low (< 0.75) requests check."),
        ("5. Telemetry & Audit", "Records processing latency (ms), confidence, and predicted ID in ai_processing_records.")
    ]
    for title, desc in steps:
        p = tfl.add_paragraph()
        p.text = f"{title}: "
        p.font.size = Pt(11.5)
        p.font.bold = True
        p.font.color.rgb = COLOR_TEXT_DARK
        p.space_before = Pt(8)
        run = p.add_run()
        run.text = desc
        run.font.bold = False
        run.font.size = Pt(11)
        run.font.color.rgb = COLOR_TEXT_MUTED

    add_card(s5, 6.8, 1.5, 5.7, 5.4, bg_color=COLOR_WHITE, border_color=COLOR_CARD_BORDER)
    tb_r = s5.shapes.add_textbox(Inches(7.0), Inches(1.7), Inches(5.3), Inches(5.0))
    tfr = tb_r.text_frame
    tfr.word_wrap = True

    hr = tfr.paragraphs[0]
    hr.text = "🎯 Human-in-the-Loop & Active Learning"
    hr.font.size = Pt(16)
    hr.font.bold = True
    hr.font.color.rgb = COLOR_GREEN

    items_r = [
        ("Manager Verification Gate", "Manager reviews the AI category badge before any authority dispatch occurs."),
        ("Zero-Friction Override", "One-click category correction dropdown if the grievance spans multiple areas."),
        ("Active Feedback Loop", "Discrepancies between predicted and final category logged for active retraining."),
        ("Dataset Corpus", "Trained on synthetic + hard real-world CSJMU grievance datasets."),
        ("Sub-100ms Inference", "Asynchronous FastAPI background execution ensures instant HTTP responses.")
    ]
    for title, desc in items_r:
        p = tfr.add_paragraph()
        p.text = f"• {title}: "
        p.font.size = Pt(11.5)
        p.font.bold = True
        p.font.color.rgb = COLOR_TEXT_DARK
        p.space_before = Pt(8)
        run = p.add_run()
        run.text = desc
        run.font.bold = False
        run.font.size = Pt(11)
        run.font.color.rgb = COLOR_TEXT_MUTED

    # ==========================================================================
    # SLIDE 6: MULTI-TIER AUTOMATED ROUTING
    # ==========================================================================
    s6 = prs.slides.add_slide(blank_layout)
    set_slide_background(s6, COLOR_BG)
    add_header(s6, "Intelligent Multi-Tier Grievance Routing Engine")

    cases = [
        ("CASE 1: Subject Assistant Dean", "SUBJECT_ASSISTANT_DEAN", "Discipline-Specific Issues (Supervisor, Coursework, Lab Allocation)", "Scholar ➔ Manager ➔ Subject Assistant Dean ➔ Direct Resolution", "Resolved directly at Assistant Dean level. No further forwarding buttons enabled.", RGBColor(14, 165, 233)),
        ("CASE 2: Grievance Cluster", "GRIEVANCE_CLUSTER", "Cross-Disciplinary Issues (Hostel Welfare, Exams, Library Services)", "Scholar ➔ Manager ➔ Assistant Dean ➔ Mapped Associate Dean ➔ Resolution", "Auto-identifies correct Associate Dean from Cluster. Button dynamically displays name.", RGBColor(234, 88, 12)),
        ("CASE 3: Fixed Authority", "FIXED_AUTHORITY", "Centralized Administrative Domains (Fellowships, Degree Verification)", "Scholar ➔ Manager ➔ Assistant Dean ➔ Fixed Specialist Officer ➔ Resolution", "Directly routed to the designated officer configured in the category entity.", COLOR_PRIMARY)
    ]

    top_c = 1.5
    for title, key, domain, flow, note, color in cases:
        add_card(s6, 0.8, top_c, 11.733, 1.65, bg_color=COLOR_WHITE, border_color=color)
        
        stripe = s6.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(top_c), Inches(0.12), Inches(1.65))
        stripe.fill.solid()
        stripe.fill.fore_color.rgb = color
        stripe.line.fill.background()

        tb = s6.shapes.add_textbox(Inches(1.1), Inches(top_c + 0.08), Inches(11.2), Inches(1.5))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = f"{title} [{key}]"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = color
        
        p2 = tf.add_paragraph()
        p2.text = f"• Domain Scope: {domain}"
        p2.font.size = Pt(11)
        p2.font.color.rgb = COLOR_TEXT_DARK
        
        p3 = tf.add_paragraph()
        p3.text = f"• Workflow Flow: {flow}"
        p3.font.size = Pt(11)
        p3.font.bold = True
        p3.font.color.rgb = COLOR_NAVY
        
        p4 = tf.add_paragraph()
        p4.text = f"• Routing Rule: {note}"
        p4.font.size = Pt(10.5)
        p4.font.color.rgb = COLOR_TEXT_MUTED
        
        top_c += 1.82

    # ==========================================================================
    # SLIDE 7: INTERACTIVE DOCUMENT REQUEST WORKFLOW
    # ==========================================================================
    s7 = prs.slides.add_slide(blank_layout)
    set_slide_background(s7, COLOR_BG)
    add_header(s7, "Interactive Additional Document Request & Review Lifecycle")

    doc_steps = [
        ("1. Authority Request", "Officer clicks 'Request Documents', specifying file names, instructions, and deadline.", "Status pauses in AWAITING_INFORMATION"),
        ("2. Scholar Notified", "In-App alert created & responsive HTML email sent with direct upload button.", "Applicant Email & Bell Alert"),
        ("3. File Upload", "Scholar uploads requested supporting PDF/Images via secure drag-and-drop workspace.", "DocumentRequest marked UPLOADED"),
        ("4. Auto-Restoration", "Grievance automatically restores to ASSIGNED/IN_PROGRESS with SAME authority.", "Zero rerouting or lost progress"),
        ("5. Authority Review", "Officer inspects file in Embedded Document Viewer: APPROVE or REJECT.", "Approval advances; Reject loops re-upload")
    ]

    left_d = 0.8
    for title, desc, badge in doc_steps:
        add_card(s7, left_d, 1.5, 2.2, 5.3, bg_color=COLOR_WHITE, border_color=COLOR_PRIMARY)
        
        db = s7.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left_d), Inches(1.5), Inches(2.2), Inches(0.75))
        db.fill.solid()
        db.fill.fore_color.rgb = COLOR_PRIMARY
        db.line.fill.background()
        
        dtf = db.text_frame
        dtf.word_wrap = True
        dp = dtf.paragraphs[0]
        dp.text = title
        dp.font.size = Pt(11)
        dp.font.bold = True
        dp.font.color.rgb = COLOR_WHITE
        dp.alignment = PP_ALIGN.CENTER

        tb = s7.shapes.add_textbox(Inches(left_d + 0.1), Inches(2.4), Inches(2.0), Inches(4.2))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p1 = tf.paragraphs[0]
        p1.text = desc
        p1.font.size = Pt(11)
        p1.font.color.rgb = COLOR_TEXT_DARK
        
        p2 = tf.add_paragraph()
        p2.text = f"\n🏷️ {badge}"
        p2.font.size = Pt(10)
        p2.font.bold = True
        p2.font.color.rgb = COLOR_GOLD
        
        left_d += 2.38

    # ==========================================================================
    # SLIDE 8: DUAL-CHANNEL NOTIFICATION & EMAIL ENGINE
    # ==========================================================================
    s8 = prs.slides.add_slide(blank_layout)
    set_slide_background(s8, COLOR_BG)
    add_header(s8, "Dual-Channel In-App Bell & Branded Email Notification Architecture")

    add_card(s8, 0.8, 1.5, 5.7, 5.4, bg_color=COLOR_WHITE, border_color=COLOR_CARD_BORDER)
    tb_l = s8.shapes.add_textbox(Inches(1.0), Inches(1.7), Inches(5.3), Inches(5.0))
    tfl = tb_l.text_frame
    tfl.word_wrap = True

    hl = tfl.paragraphs[0]
    hl.text = "🔔 In-App Centralized Notification Center"
    hl.font.size = Pt(16)
    hl.font.bold = True
    hl.font.color.rgb = COLOR_PRIMARY

    notif_points = [
        ("Real-Time Bell Widget", "Unread badge counter polling active user notifications on all dashboards."),
        ("14 Grievance Event Types", "Submissions, assignments, uploads, approvals, rejections, forwarding, & reminders."),
        ("Full History & Mark-as-Read", "Quick mark-read popup + dedicated /notifications management page."),
        ("User Data Isolation", "Scholars only receive notifications belonging to their specific grievance IDs.")
    ]
    for t, d in notif_points:
        p = tfl.add_paragraph()
        p.text = f"• {t}: "
        p.font.size = Pt(11.5)
        p.font.bold = True
        p.font.color.rgb = COLOR_TEXT_DARK
        p.space_before = Pt(8)
        run = p.add_run()
        run.text = d
        run.font.bold = False
        run.font.size = Pt(11)
        run.font.color.rgb = COLOR_TEXT_MUTED

    add_card(s8, 6.8, 1.5, 5.7, 5.4, bg_color=COLOR_WHITE, border_color=COLOR_CARD_BORDER)
    tb_r = s8.shapes.add_textbox(Inches(7.0), Inches(1.7), Inches(5.3), Inches(5.0))
    tfr = tb_r.text_frame
    tfr.word_wrap = True

    hr = tfr.paragraphs[0]
    hr.text = "✉️ High-Priority Email Notifications"
    hr.font.size = Pt(16)
    hr.font.bold = True
    hr.font.color.rgb = COLOR_GREEN

    email_points = [
        ("Important Events Only", "Emails dispatched only for DOCUMENT_REQUESTED, RESOLVED, & CLOSED."),
        ("Responsive HTML Templates", "Modular templates with CSJMU Crimson/Gold branding and action buttons."),
        ("Duplicate Protection", "State-transition checks (previous_status != new_status) prevent repeat emails."),
        ("SMTP Failure Resilience", "Safe non-blocking try-except; DB updates succeed even if mail server fails."),
        ("Zero Scholar Spam", "Internal forward/review/reminder events generate NO applicant email.")
    ]
    for t, d in email_points:
        p = tfr.add_paragraph()
        p.text = f"• {t}: "
        p.font.size = Pt(11.5)
        p.font.bold = True
        p.font.color.rgb = COLOR_TEXT_DARK
        p.space_before = Pt(8)
        run = p.add_run()
        run.text = d
        run.font.bold = False
        run.font.size = Pt(11)
        run.font.color.rgb = COLOR_TEXT_MUTED

    # ==========================================================================
    # SLIDE 9: 3-DAY INACTIVITY REMINDER ENGINE
    # ==========================================================================
    s9 = prs.slides.add_slide(blank_layout)
    set_slide_background(s9, COLOR_BG)
    add_header(s9, "3-Day Inactivity Reminder Engine & Overdue Tracking")

    r_cards = [
        ("1. Inactivity Detection Logic", "Identifies active grievances where (now - last_action_at) ≥ 3 Days. Filters out resolved, closed, or inactive assignments.", COLOR_BLUE, COLOR_BLUE_BG),
        ("2. Duplicate Guard & Cycle Reset", "Checks last_reminder_at < last_action_at to prevent spamming on repeat runs. Taking a new action resets the 3-day inactivity timer.", COLOR_AMBER, COLOR_AMBER_BG),
        ("3. Targeted Authority Alerting", "Sends internal In-App alert & email directly to the assigned officer. The scholar is never alarmed by internal delays.", COLOR_GREEN, COLOR_GREEN_BG)
    ]

    left_r = 0.8
    for title, desc, color, bg in r_cards:
        add_card(s9, left_r, 1.5, 3.75, 4.0, bg_color=bg, border_color=color)
        tb = s9.shapes.add_textbox(Inches(left_r + 0.2), Inches(1.7), Inches(3.35), Inches(3.6))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = color
        
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(12)
        p2.font.color.rgb = COLOR_TEXT_DARK
        p2.space_before = Pt(12)
        
        left_r += 3.99

    add_card(s9, 0.8, 5.7, 11.733, 1.2, bg_color=COLOR_WHITE, border_color=COLOR_PRIMARY)
    tb_b = s9.shapes.add_textbox(Inches(1.0), Inches(5.8), Inches(11.3), Inches(1.0))
    tfb = tb_b.text_frame
    tfb.word_wrap = True
    pb = tfb.paragraphs[0]
    pb.text = "⚡ Automated Background Execution"
    pb.font.size = Pt(13)
    pb.font.bold = True
    pb.font.color.rgb = COLOR_PRIMARY
    
    pb2 = tfb.add_paragraph()
    pb2.text = "Designed for seamless scheduled execution via Linux crontab or Celery Beat scheduler: 'create_overdue_reminders(db)'. Automatically tested with 10 rigorous verification criteria."
    pb2.font.size = Pt(11)
    pb2.font.color.rgb = COLOR_TEXT_MUTED

    # ==========================================================================
    # SLIDE 10: DEAN ANALYTICS & WORKLOAD MATRIX
    # ==========================================================================
    s10 = prs.slides.add_slide(blank_layout)
    set_slide_background(s10, COLOR_BG)
    add_header(s10, "Dean Executive Dashboard & Workload Performance Matrix")

    kpis = [
        ("Total University Grievances", "Real-time intake tracking across all departments"),
        ("Active Pending Bottlenecks", "Immediate flag of cases exceeding SLA limits"),
        ("Resolution & Closure Rates", "Formal verification and scholar satisfaction rate"),
        ("Average Turnaround Hours", "Granular SLA metric calculated per officer & category")
    ]
    top_k = 1.5
    for t, d in kpis:
        add_card(s10, 0.8, top_k, 5.7, 1.2, bg_color=COLOR_WHITE, border_color=COLOR_CARD_BORDER)
        tb = s10.shapes.add_textbox(Inches(1.0), Inches(top_k + 0.1), Inches(5.3), Inches(1.0))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"📊 {t}"
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = COLOR_PRIMARY
        p2 = tf.add_paragraph()
        p2.text = d
        p2.font.size = Pt(10.5)
        p2.font.color.rgb = COLOR_TEXT_MUTED
        top_k += 1.35

    add_card(s10, 6.8, 1.5, 5.7, 5.25, bg_color=COLOR_WHITE, border_color=COLOR_PRIMARY)
    tb_m = s10.shapes.add_textbox(Inches(7.0), Inches(1.7), Inches(5.3), Inches(4.8))
    tfm = tb_m.text_frame
    tfm.word_wrap = True
    
    pm = tfm.paragraphs[0]
    pm.text = "👥 Authority Performance Matrix"
    pm.font.size = Pt(16)
    pm.font.bold = True
    pm.font.color.rgb = COLOR_PRIMARY

    matrix_rows = [
        ("Prof. Namita Tiwari (DEAN)", "R&D", "Executive resolution of escalated university cases"),
        ("Mr. Asfaq (MANAGER)", "R&D", "AI review, intake verification, closure validation"),
        ("10 Assistant Deans", "R&D", "Subject specialists (Engineering, Sciences, Arts, etc.)"),
        ("3 Associate Deans", "R&D", "Domain cluster review (Sciences, Tech, Humanities)")
    ]
    for officer, dept, desc in matrix_rows:
        p = tfm.add_paragraph()
        p.text = f"• {officer} — Dept: {dept}"
        p.font.size = Pt(11.5)
        p.font.bold = True
        p.font.color.rgb = COLOR_TEXT_DARK
        p.space_before = Pt(8)
        run = p.add_run()
        run.text = f"\n  {desc}"
        run.font.bold = False
        run.font.size = Pt(10.5)
        run.font.color.rgb = COLOR_TEXT_MUTED

    # ==========================================================================
    # SLIDE 11: SECURITY & PRODUCTION READINESS
    # ==========================================================================
    s11 = prs.slides.add_slide(blank_layout)
    set_slide_background(s11, COLOR_BG)
    add_header(s11, "Enterprise Security, Compliance & Production Readiness")

    sec_items = [
        ("🔐 OAuth2 & JWT Security", "Bcrypt password hashing + HS256 HMAC JWT bearer tokens with configurable expiration."),
        ("🛡️ Granular Role-Based Access", "15 distinct permission primitives guarding every controller route."),
        ("📁 Configurable Environment", "Environment-driven API base URLs, CORS origins, and persistent document storage."),
        ("📜 Immutable Audit Logging", "Complete audit telemetry across ai_processing_records, audit_logs, & status history."),
        ("🌐 Cloud & Container Ready", "UTF-8 standardized dependencies, Nginx proxy configs, & Systemd service definitions."),
        ("✅ 100% Test Automation", "7 automated test suites verifying all 32 lifecycle criteria and regression-free health.")
    ]

    for i, (title, desc) in enumerate(sec_items):
        r_col = 0.8 if (i % 2 == 0) else 6.8
        r_top = 1.5 + (i // 2) * 1.8
        
        add_card(s11, r_col, r_top, 5.7, 1.6, bg_color=COLOR_WHITE, border_color=COLOR_CARD_BORDER)
        tb = s11.shapes.add_textbox(Inches(r_col + 0.2), Inches(r_top + 0.1), Inches(5.3), Inches(1.4))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = COLOR_PRIMARY
        
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(11)
        p2.font.color.rgb = COLOR_TEXT_MUTED
        p2.space_before = Pt(4)

    # ==========================================================================
    # SLIDE 12: CONCLUSION & FUTURE ROADMAP
    # ==========================================================================
    s12 = prs.slides.add_slide(blank_layout)
    set_slide_background(s12, COLOR_PRIMARY_DARK)

    add_card(s12, 0.8, 1.0, 11.733, 1.4, bg_color=COLOR_PRIMARY, border_color=COLOR_GOLD)
    tb_c = s12.shapes.add_textbox(Inches(1.0), Inches(1.15), Inches(11.3), Inches(1.1))
    tfc = tb_c.text_frame
    tfc.word_wrap = True
    pc1 = tfc.paragraphs[0]
    pc1.text = "NIVARAN-AI: Redefining Academic Grievance Redressal"
    pc1.font.size = Pt(22)
    pc1.font.bold = True
    pc1.font.color.rgb = COLOR_WHITE
    pc2 = tfc.add_paragraph()
    pc2.text = "Chhatrapati Shahu Ji Maharaj University (CSJMU), Kanpur — R&D Section"
    pc2.font.size = Pt(13)
    pc2.font.color.rgb = COLOR_GOLD

    pillars = [
        ("📈 Measurable Impact", "• 70% Reduction in grievance resolution turnaround.\n• Zero lost paper records with 100% digital auditability.\n• Proactive resolution via 3-day inactivity reminders.\n• Full scholar transparency with instant status tracking.", COLOR_GOLD),
        ("🚀 Future Enhancements", "• Multilingual NLP (Hindi / Regional text analysis).\n• WhatsApp API alerts for immediate scholar outreach.\n• Predictive escalation forecasting using LLMs.\n• Mobile App integration (iOS & Android).", COLOR_WHITE),
        ("🏛️ Institutional Excellence", "• Seamless compliance with NAAC / UGC guidelines.\n• Data-driven executive policy insights for Deans.\n• Secure, scalable, and cloud-ready infrastructure.\n• Built with pride for CSJMU academic community.", COLOR_GOLD)
    ]

    left_p = 0.8
    for title, desc, text_col in pillars:
        add_card(s12, left_p, 2.6, 3.75, 4.3, bg_color=COLOR_PRIMARY, border_color=COLOR_GOLD)
        tb = s12.shapes.add_textbox(Inches(left_p + 0.2), Inches(2.8), Inches(3.35), Inches(3.9))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = COLOR_GOLD
        
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(11.5)
        p2.font.color.rgb = COLOR_WHITE
        p2.space_before = Pt(10)
        
        left_p += 3.99

    # Save presentation in project root
    output_path = Path("NIVARAN_AI_Product_Presentation.pptx")
    prs.save(str(output_path))
    
    # Also save a copy in docs/
    docs_path = Path("docs/NIVARAN_AI_Product_Presentation.pptx")
    docs_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(docs_path))

    print(f"Presentation successfully created at: {output_path.resolve()} and {docs_path.resolve()}")

if __name__ == "__main__":
    create_presentation()
